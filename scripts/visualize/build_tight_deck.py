#!/usr/bin/env python3
"""Rebuild the final deck as 8 tight slides for a 5-minute video.

Each slide carries one idea, a big title, ≤5 short body lines, and a
clear visual focal point (a chart, a number callout, or a flow diagram).
The matplotlib palette of the rest of the project is reused so the deck
matches the dashboards.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# Palette — kept identical to the matplotlib charts so the deck and
# the dashboards visually rhyme.
INK = RGBColor(0x1c, 0x27, 0x36)        # near-black for body text
MUTED = RGBColor(0x55, 0x60, 0x6f)        # secondary text
SUBTLE = RGBColor(0xaa, 0xaa, 0xaa)       # tertiary
LIGHT = RGBColor(0xf3, 0xf4, 0xf6)        # surface card
PAPER = RGBColor(0xff, 0xff, 0xff)        # slide background
ORANGE = RGBColor(0xbf, 0x57, 0x00)       # primary accent (UT orange)
TEAL = RGBColor(0x00, 0x7c, 0x80)         # secondary accent
GREY_RULE = RGBColor(0xd9, 0xdc, 0xe0)    # divider lines

HEADER_FONT = "Helvetica Neue"
BODY_FONT = "Helvetica Neue"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)


def set_solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=None):
    if radius is not None:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.adjustments[0] = radius
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size, bold=False, color=INK,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=INK,
                bullet="·"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        if i > 0:
            p.space_before = Pt(8)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = BODY_FONT
    return tb


def add_title_block(slide, title, eyebrow=None):
    """Eyebrow + big title, anchored at the top of every content slide."""
    top_y = Inches(0.45)
    if eyebrow:
        add_text(slide, MARGIN, top_y, SLIDE_W - 2 * MARGIN, Inches(0.32),
                 eyebrow.upper(), size=11, bold=True, color=ORANGE,
                 font=HEADER_FONT)
        title_y = top_y + Inches(0.36)
    else:
        title_y = top_y
    add_text(slide, MARGIN, title_y, SLIDE_W - 2 * MARGIN, Inches(0.9),
             title, size=30, bold=True, color=INK, font=HEADER_FONT,
             line_spacing=1.05)
    # Hairline rule under the title
    rule = add_rect(slide, MARGIN, title_y + Inches(1.05),
                    Inches(1.1), Emu(28575), fill=ORANGE)
    return title_y + Inches(1.25)


def add_footer(slide, idx, total, label):
    add_text(slide, MARGIN, Inches(7.05),
             Inches(8.0), Inches(0.3),
             label, size=10, color=MUTED, font=BODY_FONT)
    add_text(slide, SLIDE_W - MARGIN - Inches(0.8), Inches(7.05),
             Inches(0.8), Inches(0.3),
             f"{idx} / {total}", size=10, color=MUTED, font=BODY_FONT,
             align=PP_ALIGN.RIGHT)


# ---- slide builders --------------------------------------------------


def make_title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Full-bleed orange band on the left (15% width) — the visual motif
    add_rect(s, Emu(0), Emu(0), Inches(0.35), SLIDE_H, fill=ORANGE)

    # Eyebrow
    add_text(s, Inches(0.85), Inches(0.95), Inches(11), Inches(0.32),
             "MACS 30123 · FINAL PROJECT · 2026",
             size=12, bold=True, color=ORANGE, font=HEADER_FONT)

    # Big title
    add_text(s, Inches(0.85), Inches(1.35), Inches(11.5), Inches(2.4),
             "Mapping the Discursive Battle\nOver AI on Hacker News",
             size=46, bold=True, color=INK, font=HEADER_FONT,
             line_spacing=1.05)

    # Sub
    add_text(s, Inches(0.85), Inches(3.85), Inches(11.5), Inches(0.6),
             "60 months · 411K comments · doomer vs accelerationist framings",
             size=20, color=MUTED, font=BODY_FONT)

    # Stat strip (4 boxes)
    stats = [("411K",  "comments"),
             ("60",    "months 2020–24"),
             ("4",     "stance classifiers"),
             ("4",     "aggregation backends")]
    box_w = Inches(2.7)
    gap = Inches(0.2)
    total_w = box_w * 4 + gap * 3
    x0 = (SLIDE_W - total_w) / 2
    y_box = Inches(5.05)
    for i, (big, small) in enumerate(stats):
        x = x0 + i * (box_w + gap)
        add_rect(s, x, y_box, box_w, Inches(1.55),
                 fill=LIGHT, radius=0.08)
        add_text(s, x, y_box + Inches(0.18), box_w, Inches(0.85),
                 big, size=46, bold=True, color=ORANGE,
                 font=HEADER_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x, y_box + Inches(1.0), box_w, Inches(0.4),
                 small, size=13, color=MUTED, font=BODY_FONT,
                 align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.85), Inches(6.95), Inches(11.5), Inches(0.4),
             "Cynthia Wu  ·  University of Chicago  ·  github.com/macs30123-s26/final-project-final-project",
             size=11, color=MUTED, font=BODY_FONT)


def make_question_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_y = add_title_block(
        s,
        "Did doomer vs accelerationist framing of AI\nshift around three pivotal events?",
        eyebrow="The research question",
    )

    # Left: 3 bullets
    add_bullets(s, MARGIN, body_y + Inches(0.2), Inches(7.0), Inches(3.5), [
        "Hacker News as an elite technical public sphere — narratives "
        "show up here before mainstream media.",
        "60 months · 411K comments · doomer / accelerationist / "
        "neutral stance per comment.",
        "Interrupted time series at ChatGPT (2022-11-30), with pulse "
        "dummies at GPT-4 and the OpenAI board crisis.",
    ], size=18, color=INK)

    # Right callout card
    card_x = Inches(8.05)
    card_y = body_y + Inches(0.2)
    card_w = Inches(4.7)
    card_h = Inches(3.5)
    add_rect(s, card_x, card_y, card_w, card_h, fill=LIGHT, radius=0.07)
    add_text(s, card_x + Inches(0.4), card_y + Inches(0.35),
             card_w - Inches(0.8), Inches(0.4),
             "THE METHODOLOGICAL POINT", size=11, bold=True,
             color=ORANGE, font=HEADER_FONT)
    add_text(s, card_x + Inches(0.4), card_y + Inches(0.85),
             card_w - Inches(0.8), card_h - Inches(1.2),
             "Single-classifier discourse analysis is fragile. "
             "A defensible answer requires running the same panel "
             "under multiple independent measurement instruments and "
             "seeing where they agree or disagree.",
             size=16, color=INK, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, 2, 8, "Research question")


def make_why_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_y = add_title_block(
        s,
        "Why this is a scalable-computing problem",
        eyebrow="Scale, three ways",
    )

    cards = [
        ("Corpus scale",
         "1,260 partitioned shards across 60 months · "
         "year × keyword family × query",
         "1,260"),
        ("Per-doc classification × 4",
         "411K comments scored by four independent stance instruments — "
         "CPU + GPU.",
         "411K × 4"),
        ("Aggregation × 4 backends",
         "Same 1,260 inference parquets reduced under serial pandas, "
         "Dask local, Dask multi-node, PySpark.",
         "60-month panel"),
    ]
    card_w = (SLIDE_W - 2 * MARGIN - Inches(0.5)) / 3
    gap = Inches(0.25)
    card_h = Inches(4.0)
    for i, (head, body, big) in enumerate(cards):
        x = MARGIN + i * (card_w + gap)
        add_rect(s, x, body_y + Inches(0.2), card_w, card_h,
                 fill=PAPER, line=GREY_RULE,
                 line_w=Emu(9525), radius=0.05)
        # Number band
        add_rect(s, x, body_y + Inches(0.2), card_w, Inches(0.18),
                 fill=ORANGE)
        # Card content
        add_text(s, x + Inches(0.35), body_y + Inches(0.55),
                 card_w - Inches(0.7), Inches(0.45),
                 head, size=18, bold=True, color=INK,
                 font=HEADER_FONT)
        add_text(s, x + Inches(0.35), body_y + Inches(1.15),
                 card_w - Inches(0.7), Inches(1.0),
                 big, size=34, bold=True, color=ORANGE,
                 font=HEADER_FONT)
        add_text(s, x + Inches(0.35), body_y + Inches(2.2),
                 card_w - Inches(0.7), Inches(1.7),
                 body, size=15, color=MUTED, font=BODY_FONT,
                 line_spacing=1.3)

    add_footer(s, 3, 8, "Why scalable computing")


def make_pipeline_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_y = add_title_block(
        s,
        "Pipeline — 6 stages, partitioned files, re-run any branch",
        eyebrow="Architecture",
    )

    stages = [
        ("①", "Collect",  "AWS EC2 + S3\n+ local parallel",  ORANGE),
        ("②", "Clean + 4 classifiers", "Midway Slurm\nCPU + GPU array", TEAL),
        ("③", "Aggregate × 4 backends", "Midway\npandas / Dask / Spark", TEAL),
        ("④", "ITS + atlas + topics", "Midway Slurm\nstatsmodels / BERTopic", TEAL),
        ("⑤", "Gold-set validation", "Midway Slurm\n200-row hand audit", TEAL),
        ("⑥", "Public dashboards", "AWS S3\nstatic website", ORANGE),
    ]
    # 3 × 2 grid
    cols = 3
    rows = 2
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)
    card_w = (SLIDE_W - 2 * MARGIN - gap_x * (cols - 1)) / cols
    card_h = Inches(2.05)
    grid_y = body_y + Inches(0.2)
    for i, (num, head, body, color) in enumerate(stages):
        r, c = i // cols, i % cols
        x = MARGIN + c * (card_w + gap_x)
        y = grid_y + r * (card_h + gap_y)
        add_rect(s, x, y, card_w, card_h, fill=LIGHT, radius=0.05)
        # Big circled number
        add_text(s, x + Inches(0.25), y + Inches(0.1),
                 Inches(0.6), Inches(0.6),
                 num, size=28, bold=True, color=color, font=HEADER_FONT)
        # Stage head
        add_text(s, x + Inches(0.95), y + Inches(0.18),
                 card_w - Inches(1.2), Inches(0.5),
                 head, size=18, bold=True, color=INK, font=HEADER_FONT)
        # Stage body
        add_text(s, x + Inches(0.95), y + Inches(0.78),
                 card_w - Inches(1.2), card_h - Inches(0.9),
                 body, size=13, color=MUTED, font=BODY_FONT,
                 line_spacing=1.3)

    add_text(s, MARGIN, Inches(6.6), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "Everything after Stage ① runs on Midway Slurm — CPU arrays, GPU, Dask multi-node, PySpark.",
             size=12, color=MUTED, font=BODY_FONT, align=PP_ALIGN.LEFT)

    add_footer(s, 4, 8, "Pipeline")


def make_benchmark_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_y = add_title_block(
        s,
        "Four aggregation backends, same 1,260 parquets",
        eyebrow="Backend benchmark",
    )

    img_path = Path("deliverables/framework_benchmark.png")
    if img_path.exists():
        # Center the image, leave room for caption.
        max_w = Inches(11.0)
        pic = s.shapes.add_picture(str(img_path),
                                   Inches(0), body_y + Inches(0.2),
                                   width=max_w)
        max_h = Inches(4.3)
        if pic.height > max_h:
            scale = max_h / pic.height
            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)
        pic.left = int((SLIDE_W - pic.width) / 2)

    # Punchline strip
    add_rect(s, MARGIN, Inches(6.0), SLIDE_W - 2 * MARGIN, Inches(0.8),
             fill=LIGHT, radius=0.05)
    add_text(s, MARGIN + Inches(0.3), Inches(6.12),
             SLIDE_W - 2 * MARGIN - Inches(0.6), Inches(0.7),
             "PySpark 23 s — Dask multi-node 36 s — pandas 62 s — Dask local 152 s. "
             "Distributed isn’t automatically faster.",
             size=15, bold=True, color=INK, font=BODY_FONT)

    add_footer(s, 5, 8, "Backend benchmark")


def make_its_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_y = add_title_block(
        s,
        "7 of 8 headline coefficients disagree across 4 classifiers",
        eyebrow="Cross-method ITS",
    )

    # Big stat on left
    add_text(s, MARGIN, body_y + Inches(0.3),
             Inches(4.6), Inches(2.0),
             "7 / 8", size=110, bold=True, color=ORANGE,
             font=HEADER_FONT, line_spacing=1.0)
    add_text(s, MARGIN, body_y + Inches(2.4),
             Inches(4.6), Inches(0.7),
             "headline coefficients\nfail the agreement criterion",
             size=14, color=MUTED, font=BODY_FONT, line_spacing=1.3)

    # Bottom callout: 3 sign flips
    add_rect(s, MARGIN, body_y + Inches(3.4),
             Inches(4.6), Inches(1.4),
             fill=LIGHT, radius=0.05)
    add_text(s, MARGIN + Inches(0.25), body_y + Inches(3.55),
             Inches(4.2), Inches(0.4),
             "3 SIGN FLIPS", size=11, bold=True, color=ORANGE,
             font=HEADER_FONT)
    add_text(s, MARGIN + Inches(0.25), body_y + Inches(3.95),
             Inches(4.2), Inches(1.0),
             "doomer × ChatGPT level · accel × board crisis · accel × GPT-4",
             size=14, color=INK, font=BODY_FONT, line_spacing=1.3)

    # Right: surviving direction
    card_x = Inches(6.4)
    card_w = SLIDE_W - card_x - MARGIN
    card_h = Inches(4.8)
    add_rect(s, card_x, body_y + Inches(0.3), card_w, card_h,
             fill=TEAL, radius=0.05)
    add_text(s, card_x + Inches(0.35), body_y + Inches(0.55),
             card_w - Inches(0.7), Inches(0.45),
             "ONE DIRECTION SURVIVES EVERY METHOD",
             size=12, bold=True, color=PAPER, font=HEADER_FONT)
    add_text(s, card_x + Inches(0.35), body_y + Inches(1.15),
             card_w - Inches(0.7), Inches(1.6),
             "Board-crisis doomer pulse",
             size=28, bold=True, color=PAPER, font=HEADER_FONT,
             line_spacing=1.1)
    add_text(s, card_x + Inches(0.35), body_y + Inches(2.85),
             card_w - Inches(0.7), Inches(1.8),
             "+0.012 to +0.131 across TF-IDF, zero-shot gated, "
             "and zero-shot ungated — every method finds a rise "
             "in doomer share after the OpenAI board episode.",
             size=15, color=PAPER, font=BODY_FONT, line_spacing=1.3)

    add_footer(s, 6, 8, "Cross-method ITS")


def make_gold_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    body_y = add_title_block(
        s,
        "Gold-set vs human judgement — bounding the instrument",
        eyebrow="Validation",
    )

    img_path = Path("deliverables/confusion_matrices.png")
    if img_path.exists():
        max_w = Inches(12.2)
        pic = s.shapes.add_picture(str(img_path),
                                   Inches(0), body_y + Inches(0.05),
                                   width=max_w)
        max_h = Inches(3.4)
        if pic.height > max_h:
            scale = max_h / pic.height
            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)
        pic.left = int((SLIDE_W - pic.width) / 2)

    # Two stat callouts under the image
    cards = [
        ("51 %",  "lexicon precision floor",
         "of 200 hand-labelled rows are wrong_other — retrieved by the "
         "lexicon but not actually about AI."),
        ("98.3 %", "embedding recall ceiling",
         "MiniLM cosine-similarity cross-check — the lexicon "
         "over-retrieves but rarely under-retrieves."),
    ]
    card_y = Inches(5.45)
    card_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    for i, (big, head, body) in enumerate(cards):
        x = MARGIN + i * (card_w + Inches(0.3))
        add_rect(s, x, card_y, card_w, Inches(1.45),
                 fill=LIGHT, radius=0.05)
        add_text(s, x + Inches(0.3), card_y + Inches(0.15),
                 Inches(2.2), Inches(0.95),
                 big, size=42, bold=True, color=ORANGE,
                 font=HEADER_FONT, line_spacing=1.0)
        add_text(s, x + Inches(2.55), card_y + Inches(0.2),
                 card_w - Inches(2.8), Inches(0.4),
                 head.upper(), size=11, bold=True, color=ORANGE,
                 font=HEADER_FONT)
        add_text(s, x + Inches(2.55), card_y + Inches(0.6),
                 card_w - Inches(2.8), Inches(1.0),
                 body, size=12, color=INK, font=BODY_FONT,
                 line_spacing=1.3)

    add_footer(s, 7, 8, "Gold-set validation")



def make_close_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark hero banner across the top half
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(3.4), fill=INK)
    add_text(s, MARGIN, Inches(0.7), SLIDE_W - 2 * MARGIN, Inches(0.4),
             "WHAT THE PIPELINE BUYS YOU",
             size=12, bold=True, color=ORANGE, font=HEADER_FONT)
    add_text(s, MARGIN, Inches(1.2), SLIDE_W - 2 * MARGIN, Inches(1.8),
             "A 60-month panel of HN AI discourse,\n"
             "cross-checked under multiple instruments.",
             size=34, bold=True, color=PAPER, font=HEADER_FONT,
             line_spacing=1.1)

    # Three callouts in the lower half
    items = [
        ("4 stance classifiers",
         "lexicon · TF-IDF + LR · zero-shot DeBERTa · DistilBERT fine-tune"),
        ("4 aggregation backends",
         "pandas · Dask local · Dask multi-node · PySpark"),
        ("3 collection venues",
         "Local parallel · AWS EC2 + S3 · Google BigQuery"),
    ]
    body_y = Inches(3.9)
    card_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 3
    for i, (head, body) in enumerate(items):
        x = MARGIN + i * (card_w + Inches(0.2))
        add_rect(s, x, body_y, card_w, Inches(2.0),
                 fill=PAPER, line=GREY_RULE, line_w=Emu(9525), radius=0.05)
        add_rect(s, x, body_y, card_w, Inches(0.15), fill=ORANGE)
        add_text(s, x + Inches(0.3), body_y + Inches(0.4),
                 card_w - Inches(0.6), Inches(0.55),
                 head, size=18, bold=True, color=INK, font=HEADER_FONT)
        add_text(s, x + Inches(0.3), body_y + Inches(1.05),
                 card_w - Inches(0.6), Inches(0.9),
                 body, size=13, color=MUTED, font=BODY_FONT,
                 line_spacing=1.3)

    # Live dashboard URL line
    add_rect(s, MARGIN, Inches(6.25), SLIDE_W - 2 * MARGIN,
             Inches(0.7), fill=LIGHT, radius=0.05)
    add_text(s, MARGIN + Inches(0.3), Inches(6.32),
             Inches(3.5), Inches(0.55),
             "LIVE DASHBOARDS", size=11, bold=True, color=ORANGE,
             font=HEADER_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MARGIN + Inches(3.5), Inches(6.32),
             SLIDE_W - 2 * MARGIN - Inches(3.6), Inches(0.55),
             "hn-ai-discourse-public-2026.s3-website-us-east-1.amazonaws.com",
             size=14, color=INK, font=BODY_FONT,
             anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 8, 8, "Contribution")


# ---- main -----------------------------------------------------------


def main():
    out = Path("deliverables/hn-ai-discourse-final-presentation-2020-2024.pptx")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)
    make_question_slide(prs)
    make_why_slide(prs)
    make_pipeline_slide(prs)
    make_benchmark_slide(prs)
    make_its_slide(prs)
    make_gold_slide(prs)
    make_close_slide(prs)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    size_kb = out.stat().st_size / 1024
    print(f"wrote {out}  ({len(prs.slides)} slides, {size_kb:.0f} KB)")


if __name__ == "__main__":
    main()
